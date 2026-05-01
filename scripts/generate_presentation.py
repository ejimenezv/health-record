"""Generate the MedRecord presentation deck (.pptx) from the script."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "presentation" / "medrecord-presentation.pptx"

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
TEAL = RGBColor(0x12, 0x7B, 0x7B)
SLATE = RGBColor(0x33, 0x33, 0x33)
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
ACCENT = RGBColor(0xE0, 0x6C, 0x00)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
GREY = RGBColor(0x77, 0x77, 0x77)


def add_background(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text, *, size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def add_bullets(slide, left, top, width, height, items, *, size=18, color=SLATE):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def add_accent_bar(slide, color=TEAL):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.7), prs.slide_width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_footer(slide, page_num):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
             "MedRecord  ·  AI/LLM Solution Architect  ·  Cohorte 2025-A",
             size=10, color=GREY)
    add_text(slide, Inches(12), Inches(7.1), Inches(1), Inches(0.3),
             str(page_num), size=10, color=GREY, align=PP_ALIGN.RIGHT)


def title_slide(prs, title, subtitle, author):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_background(slide, NAVY)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.4), Inches(0.15), Inches(2.6))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    add_text(slide, Inches(0.9), Inches(2.4), Inches(12), Inches(1.2), title,
             size=54, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(slide, Inches(0.9), Inches(3.6), Inches(12), Inches(1.0), subtitle,
             size=22, color=RGBColor(0xCF, 0xD8, 0xDC))
    add_text(slide, Inches(0.9), Inches(4.5), Inches(12), Inches(0.6), author,
             size=18, color=RGBColor(0xCF, 0xD8, 0xDC))
    add_text(slide, Inches(0.9), Inches(6.6), Inches(12), Inches(0.5),
             "Proyecto Final  ·  AI/LLM Solution Architect  ·  Cohorte 2025-A",
             size=14, color=ACCENT)


def section_slide(prs, title, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, LIGHT)
    add_accent_bar(slide, ACCENT)
    add_text(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.5),
             "MedRecord", size=14, bold=True, color=NAVY)
    add_text(slide, Inches(0.6), Inches(3.0), Inches(12), Inches(2),
             title, size=44, bold=True, color=NAVY)
    add_footer(slide, page_num)
    return slide


def content_slide(prs, title, page_num, accent=TEAL):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, RGBColor(0xFF, 0xFF, 0xFF))
    add_accent_bar(slide, accent)
    add_text(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.5),
             "MedRecord", size=14, bold=True, color=NAVY)
    add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    add_footer(slide, page_num)
    return slide


def add_table(slide, left, top, width, height, headers, rows, *, header_color=NAVY, body_color=SLATE):
    cols = len(headers)
    tbl = slide.shapes.add_table(len(rows) + 1, cols, left, top, width, height).table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = h
        r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        r.font.name = "Calibri"
        cell.fill.solid(); cell.fill.fore_color.rgb = header_color
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(12); r.font.color.rgb = body_color
            r.font.name = "Calibri"
            if i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return tbl


# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

page = 1

# 1. Cover
title_slide(prs,
            "MedRecord",
            "Sistema de documentación médica automatizada en español",
            "Enrique Jiménez")
page += 1

# 2. El problema
slide = content_slide(prs, "El problema", page)
add_text(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(0.6),
         "30–40 % del tiempo del médico se va en documentación post-consulta.",
         size=22, bold=True, color=NAVY)
add_bullets(slide, Inches(0.6), Inches(2.8), Inches(12), Inches(3.5), [
    "Notas SOAP redactadas a mano, horas después de la consulta.",
    "Información perdida o incompleta entre la consulta y el expediente.",
    "Errores de transcripción difíciles de auditar.",
    "Alertas de interacción medicamentosa que llegan tarde — o no llegan.",
], size=20)
page += 1

# 3. La propuesta
slide = content_slide(prs, "La propuesta", page)
add_text(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(0.7),
         "Transcripción y extracción clínica en vivo, mientras la consulta ocurre.",
         size=22, bold=True, color=NAVY)
add_bullets(slide, Inches(0.6), Inches(2.8), Inches(12), Inches(3.5), [
    "Streaming bidireccional con WebSocket desde el inicio.",
    "Diarización doctor / paciente híbrida (acústica + LLM).",
    "Extracción incremental de síntomas, diagnósticos, CIE-10, fármacos.",
    "Validación contra base de conocimiento médico vía RAG.",
    "Alertas de interacción medicamentosa emitidas en cuanto se detectan.",
], size=20)
page += 1

# 4. Stack
slide = content_slide(prs, "Stack tecnológico", page)
add_table(slide, Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.5),
          ["Capa", "Tecnología", "Rol"],
          [
              ["Frontend", "React + TypeScript + Vite", "SPA con grabación en vivo y panel de transcripción"],
              ["Backend de negocio", "Node.js + Express + Prisma", "Pacientes, citas, expediente, JWT"],
              ["Servicio de IA", "Python + FastAPI", "Streaming, transcripción, extracción, RAG"],
              ["Vector store", "ChromaDB + text-embedding-3-small", "Vademécum y guías clínicas"],
              ["Datos", "PostgreSQL + Redis", "Metadatos clínicos / sesión / cache"],
              ["Modelos", "OpenAI Whisper + GPT-4o family", "Audio y extracción"],
              ["Infraestructura", "Docker Compose · AWS · Terraform", "Deployment reproducible"],
          ])
page += 1

# 5. Arquitectura (placeholder — usuario reemplaza con C4)
slide = content_slide(prs, "Arquitectura — vista de contenedores", page)
add_text(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(0.5),
         "Microservicios con frontera clínica / IA bien separada.",
         size=18, color=SLATE)
add_bullets(slide, Inches(0.6), Inches(2.4), Inches(12), Inches(4.0), [
    "Frontend (React) ⇄ REST + Socket.IO ⇄ Backend (Node) ⇄ AI Service (Python).",
    "Backend de Node es deliberadamente delgado en lo que toca a IA: invoca, no orquesta.",
    "El servicio de IA concentra streaming, Whisper, diarización, extracción, RAG y cost tracking.",
    "Datos: PostgreSQL para clínico, Redis para sesión y cache, ChromaDB para vectores.",
    "Documentación visual: docs/architecture/diagrams/  (sustituir esta diapositiva por el diagrama C4).",
], size=18)
page += 1

# 6. Demo overview
slide = content_slide(prs, "Demo en vivo — flujo", page, accent=ACCENT)
add_table(slide, Inches(0.6), Inches(1.6), Inches(12.1), Inches(4.8),
          ["#", "Paso", "Qué se muestra"],
          [
              ["1", "Health checks",       "Backend y AI service responden 'healthy'."],
              ["2", "Login + dashboard",   "JWT compartido; navegación clínica."],
              ["3", "Paciente y cita",     "CRUD; alergias y crónicas como entidades."],
              ["4", "Sesión en vivo",      "WebSocket; transcripción incremental."],
              ["5", "Diarización",         "Turnos doctor / paciente en tiempo real."],
              ["6", "Extracción en vivo",  "Síntomas, diagnóstico, CIE-10 candidato."],
              ["7", "Alerta de interacción", "Drug-interaction warning vía WebSocket."],
              ["8", "Cierre de sesión",    "Transcripción + extracción consolidadas."],
              ["9", "Consulta RAG",        "POST /query con citaciones y similitud."],
          ])
page += 1

# 7. ADR-001
slide = content_slide(prs, "ADR-001 · Estrategia multi-tier de LLMs", page)
add_text(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(0.5),
         "Tres tiers para optimizar costo sin sacrificar calidad clínica.",
         size=18, bold=True, color=TEAL)
add_table(slide, Inches(0.6), Inches(2.3), Inches(12.1), Inches(2.4),
          ["Tier", "Modelo", "Uso"],
          [
              ["FAST_CHEAP", "GPT-4o-mini", "Validaciones rápidas, validador de diarización, formato"],
              ["BALANCED",   "GPT-4o",      "Extracción clínica, generación principal — baseline"],
              ["PREMIUM",    "GPT-4-turbo", "Fallback para casos complejos / JSON inválido"],
          ])
add_text(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(1.5),
         "Trade-off: GPT-4o-mini es ~16× más barato pero menos preciso en dosis. "
         "GPT-4o queda como baseline; mini se usa solo en validaciones secundarias. "
         "PREMIUM solo si BALANCED falla.",
         size=15, color=SLATE)
page += 1

# 8. ADR-002
slide = content_slide(prs, "ADR-002 · ChromaDB sobre Pinecone", page)
add_table(slide, Inches(0.6), Inches(1.6), Inches(12.1), Inches(2.0),
          ["Opción", "Costo / mes", "Latencia retrieval", "Veredicto"],
          [
              ["ChromaDB (elegido)", "$0 (local)",   "80–120 ms", "Suficiente al volumen actual"],
              ["Pinecone",           "~$50 starter", "50–80 ms",  "Sin replicación / backup ganados no justifican costo"],
              ["Weaviate",           "~$25 cloud",  "60–100 ms", "Complejidad operacional alta"],
          ])
add_text(slide, Inches(0.6), Inches(4.0), Inches(12), Inches(2.5),
         "El cuello de botella real es Whisper, no el retrieval. Optimizar el componente "
         "que pesa el 0,1 % del tiempo no tiene sentido. Criterios de revisión: > 80 k vectores, "
         "retrieval p95 > 500 ms o necesidad multi-región.",
         size=15, color=SLATE)
page += 1

# 9. ADR-003
slide = content_slide(prs, "ADR-003 · Node + Python con JWT compartido", page)
add_bullets(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(4.5), [
    "Node.js para dominio clínico: ecosistema Prisma + Express maduro y conocido.",
    "Python para IA: Whisper, RAGAS, ChromaDB, Resemblyzer, sin equivalente real en Node.",
    "JWT HS256 con secreto compartido byte-a-byte: un solo token vale en ambos servicios.",
    "Frontera clara: el backend de Node nunca toca prompts; el servicio de Python nunca toca pacientes.",
    "Permite escalar y desplegar cada servicio en aislamiento.",
], size=18)
page += 1

# 10. ADR-005
slide = content_slide(prs, "ADR-005 · Diarización híbrida", page)
add_text(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(0.6),
         "AudioFeatureDiarizer + LLMValidator + IncrementalDiarizer.",
         size=18, bold=True, color=TEAL)
add_table(slide, Inches(0.6), Inches(2.4), Inches(12.1), Inches(1.8),
          ["Modo", "Precisión", "Latencia adicional"],
          [
              ["Streaming (online)",         "~ 87 %", "< 2 s por chunk"],
              ["Refinamiento batch al cerrar","~ 92 %", "+ 30 s al cerrar la sesión"],
          ])
add_text(slide, Inches(0.6), Inches(4.6), Inches(12), Inches(2.0),
         "Pyannote (estado del arte) descartado: requiere GPU; en CPU es incompatible "
         "con tiempo real. Resemblyzer descartado: ~80 % de precisión, no mejora el baseline. "
         "El híbrido alcanza precisión clínicamente útil sin GPU.",
         size=15, color=SLATE)
page += 1

# 11. ADR-006
slide = content_slide(prs, "ADR-006 · Streaming en tiempo real", page)
add_bullets(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(2.5), [
    "WebSocket bidireccional desde el día uno (no agregado después).",
    "Polling y SSE evaluados y descartados (latencia / unidireccionalidad).",
    "VAD del lado servidor (Silero) decide cuándo enviar audio a Whisper.",
], size=18)
add_text(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(0.5),
         "Intelligent buffering",
         size=18, bold=True, color=TEAL)
add_bullets(slide, Inches(0.6), Inches(4.9), Inches(12), Inches(2.0), [
    "Voz activa → chunks cada 5 s (priorizar latencia).",
    "Silencio 2–10 s → batch (optimizar costo).",
    "Silencio prolongado → no se envía nada.",
], size=16)
page += 1

# 12. Resultados — RAGAS
slide = content_slide(prs, "Resultados · RAGAS (medido)", page, accent=GREEN)
add_text(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(0.6),
         "Ejecución 2026-04-30 · 8 preguntas sintéticas validadas",
         size=14, color=GREY)
add_table(slide, Inches(0.6), Inches(2.3), Inches(12.1), Inches(2.4),
          ["Métrica", "Score", "Umbral", "Estado"],
          [
              ["Faithfulness",     "0.938", "> 0.80", "Pass"],
              ["Context Precision","1.000", "> 0.75", "Pass"],
              ["Answer Relevancy", "0.964", "> 0.75", "Pass"],
              ["Context Recall",   "1.000", "> 0.70", "Pass"],
          ])
add_text(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(1.8),
         "Caveat explícito: dataset sintético de 8 preguntas. Es un guardrail de regresión, "
         "no una prueba de calidad clínica en producción. Evaluación con consultas reales "
         "es trabajo de v2.",
         size=15, color=SLATE)
page += 1

# 13. Resultados — Carga
slide = content_slide(prs, "Resultados · Carga y latencia (medido)", page, accent=GREEN)
add_table(slide, Inches(0.6), Inches(1.6), Inches(12.1), Inches(3.0),
          ["Escenario", "Métrica", "Valor", "Objetivo", "Estado"],
          [
              ["Persistencia de eventos",       "Write p95",         "14.45 ms",     "< 50 ms",  "Pass"],
              ["Persistencia de eventos",       "Throughput",        "712 writes/s", "≥ 50 w/s", "Pass"],
              ["Persistencia de eventos",       "Error rate",        "0.00 %",       "< 1 %",    "Pass"],
              ["WebSocket handshake",           "Latencia mediana",  "59 ms",        "< 500 ms", "Pass"],
          ])
add_text(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(1.8),
         "Fuente: ai-service/reports/2026-04-30/load_test_report.md.  "
         "Throughput sostenido de streaming end-to-end queda fuera del alcance de esta "
         "tabla y se reporta cualitativamente en la demo.",
         size=14, color=GREY)
page += 1

# 14. Resultados — Costos
slide = content_slide(prs, "Resultados · Costos", page, accent=ACCENT)
add_bullets(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(3.0), [
    "Cost tracker instrumentado en cada llamada a OpenAI.",
    "Endpoint backend GET /api/v1/costs con desglose por servicio y modo.",
    "Caching de embeddings en Redis para evitar llamadas redundantes.",
    "Modelo de costos proyectado: $0.25–0.28 por consulta en tiempo real.",
], size=18)
add_text(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(1.8),
         "Honestidad explícita: la cifra $0.25–0.28 es del modelo de costos en "
         "docs/delivery-4/02-cost-analysis.md, no un promedio observado sobre sesiones "
         "completadas. Reconciliación con facturación AWS = OI-5, prioridad de v2.",
         size=14, color=GREY)
page += 1

# 15. Lo que funcionó
slide = content_slide(prs, "Reflexión · Lo que funcionó", page, accent=GREEN)
add_bullets(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(4.5), [
    "Multi-tier LLM (ADR-001): GPT-4o donde importa, mini para validaciones secundarias.",
    "ChromaDB local (ADR-002): cero costo operacional, latencia suficiente al volumen actual.",
    "Diarización híbrida (ADR-005): 87 % streaming, 92 % batch, sin GPU.",
    "Streaming desde día uno (ADR-006): habilitó alertas en cuanto se detectan.",
    "RAGAS como guardrail de regresión integrado en el ciclo de desarrollo.",
], size=18, color=SLATE)
page += 1

# 16. Lo que no funcionó
slide = content_slide(prs, "Reflexión · Lo que quedó pendiente", page, accent=RED)
add_bullets(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(4.5), [
    "Generación automática del borrador SOAP — modelo de datos listo, orquestación pendiente.",
    "Dependencia de la API de Whisper para latencia y costo variable.",
    "Sin dashboard de costos en el frontend (endpoint backend sí existe).",
    "Sin Grafana / Prometheus — solo logs estructurados y métricas in-process.",
    "Reconciliación con facturación real de AWS pendiente (OI-5).",
], size=18, color=SLATE)
page += 1

# 17. Lecciones
slide = content_slide(prs, "Lecciones aprendidas", page, accent=TEAL)
add_bullets(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(4.5), [
    "Prompt engineering + few-shot vence a fine-tuning en la fase MVP.",
    "Cost tracking es un requisito no funcional, no una feature: se instrumenta primero.",
    "RAGAS detecta problemas de fidelidad invisibles a los tests unitarios.",
    "Streaming en tiempo real se diseña desde el inicio, no se retrofitea.",
    "Honestidad sobre lo medido vs. lo modelado importa más que pintar todo verde.",
], size=18)
page += 1

# 18. Roadmap
slide = content_slide(prs, "Roadmap v2", page, accent=ACCENT)
add_text(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(0.5),
         "Prioridad alta", size=18, bold=True, color=NAVY)
add_bullets(slide, Inches(0.6), Inches(2.2), Inches(12), Inches(2.5), [
    "Generación automática del borrador SOAP a partir de la extracción consolidada.",
    "Migración a faster-whisper local (latencia + eliminación de costo variable).",
    "Dashboard de costos en frontend + observabilidad con Langfuse y Grafana.",
], size=16)
add_text(slide, Inches(0.6), Inches(4.7), Inches(12), Inches(0.5),
         "Medio / largo plazo", size=18, bold=True, color=NAVY)
add_bullets(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(2.0), [
    "Multi-tenancy con autenticación federada (Azure AD / Cognito).",
    "Soporte multimodal con GPT-4-vision para imágenes clínicas.",
    "Alertas más ricas basadas en historial completo del paciente.",
], size=16)
page += 1

# 19. Cierre
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_text(slide, Inches(0.9), Inches(2.4), Inches(12), Inches(1.2),
         "Gracias", size=66, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(slide, Inches(0.9), Inches(3.7), Inches(12), Inches(0.6),
         "MedRecord — documentación médica automatizada en español",
         size=20, color=RGBColor(0xCF, 0xD8, 0xDC))
add_text(slide, Inches(0.9), Inches(4.7), Inches(12), Inches(0.5),
         "Repositorio:  github.com/<tu-usuario>/health-record",
         size=16, color=ACCENT)
add_text(slide, Inches(0.9), Inches(5.2), Inches(12), Inches(0.5),
         "Documentación:  docs/delivery-4/  ·  docs/adr/",
         size=16, color=ACCENT)
add_text(slide, Inches(0.9), Inches(6.6), Inches(12), Inches(0.5),
         "Enrique Jiménez  ·  AI/LLM Solution Architect  ·  Cohorte 2025-A",
         size=14, color=RGBColor(0xCF, 0xD8, 0xDC))

OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUTPUT)
print(f"Wrote {OUTPUT}  ·  {len(prs.slides)} slides")
